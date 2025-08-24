# pptgen_module.py

# pptgen_module.py

import os
import re
from datetime import datetime
from typing import List, Optional, Iterable

import numpy as np
import pandas as pd
from dateutil.relativedelta import relativedelta
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


_QUARTER_PATTERNS = [
    re.compile(r"^\s*Q([1-4])[-\s]?(\d{4})\s*$", re.IGNORECASE),   # Q1-2025 or Q1 2025
    re.compile(r"^\s*(\d{4})[-\s]?Q([1-4])\s*$", re.IGNORECASE),   # 2025-Q1 or 2025 Q1
    re.compile(r"^\s*([1-4])\s*[Qq][-]?\s*(\d{4})\s*$"),           # 1Q-2025 / 1q 2025
]


def _quarter_to_period(s: str) -> Optional[pd.Period]:
    """Parse various quarter strings to a pandas Period (Q-DEC). Expecting formats like 'Q1-2025'."""
    if s is None or (isinstance(s, float) and np.isnan(s)):
        return None
    text = str(s).strip()
    for pat in _QUARTER_PATTERNS:
        m = pat.match(text)
        if m:
            if pat.pattern.startswith("^\\s*Q"):
                q = int(m.group(1)); year = int(m.group(2))
            elif pat.pattern.startswith("^\\s*(\\d{4})"):
                year = int(m.group(1)); q = int(m.group(2))
            else:
                q = int(m.group(1)); year = int(m.group(2))
            # pandas uses Q-DEC by default; Period(f"{year}Q{q}") is fine.
            try:
                return pd.Period(f"{year}Q{q}", freq="Q")
            except Exception:
                return None
    # Last resort: accept already like '2025Q1'
    try:
        return pd.Period(str(text).replace(" ", ""), freq="Q")
    except Exception:
        return None


def _period_to_label(p: pd.Period) -> str:
    """Return canonical display like '2025Q1'."""
    return f"{p.year}Q{p.quarter}"


class MonthlyPerformancePPT:
    """
    Generate per-country PowerPoints from *quarterly* performance data.

    Provide a column_map to specify which columns to use:
      {
        "country": "<exact col name>",     # required
        "model":   "<exact col name>",     # required
        "quarter": "<exact col name>",     # required, values like 'Q1-2025'
        "units":   "<exact col name>",     # required
        "revenue": "<exact col name or None>",  # optional
        "price":   "<exact col name or None>"   # optional
      }
    """

    def __init__(
        self,
        output_dir: str = "ppt_output",
        template_ppt: Optional[str] = None,
        logo_path: Optional[str] = None,
        last_n_quarters: int = 6,
    ):
        self.output_dir = output_dir
        self.template_ppt = template_ppt    # e.g., "BAL.pptx"
        self.logo_path = logo_path          # e.g., "company_logo.png"
        self.last_n_quarters = max(1, int(last_n_quarters))

        self._ensure_dir(self.output_dir)
        self._tmp_dir = os.path.join(self.output_dir, "_tmp")
        self._ensure_dir(self._tmp_dir)

    # ---------- Public API ----------

    def generate_from_file(
        self,
        input_path: str,
        sheet_name: Optional[str] = None,
        column_map: dict | None = None,
    ) -> List[str]:
        df = self._read_table(input_path, sheet_name)
        work = self._prepare_dataframe(df, column_map=column_map)
        return self._generate_all(work)

    def generate_from_dataframe(self, df: pd.DataFrame, column_map: dict | None = None) -> List[str]:
        work = self._prepare_dataframe(df, column_map=column_map)
        return self._generate_all(work)

    # ---------- Core flow ----------

    def _generate_all(self, work: pd.DataFrame) -> List[str]:
        if work.empty:
            raise ValueError("No rows after preprocessing. Check your input and column mapping.")

        latest_q = work["QuarterPeriod"].max()  # pandas Period
        start_q = latest_q - (self.last_n_quarters - 1)

        outputs: List[str] = []
        for country, df_c in work.groupby("Country"):
            out_path = self._generate_for_country(
                country=country,
                df_country=df_c.copy(),
                latest_q=latest_q,
                start_q=start_q,
            )
            outputs.append(out_path)

        return outputs

    def _generate_for_country(
        self,
        country: str,
        df_country: pd.DataFrame,
        latest_q: pd.Period,
        start_q: pd.Period,
    ) -> str:
        prs = Presentation(self.template_ppt) if (self.template_ppt and os.path.exists(self.template_ppt)) else Presentation()

        # ---- Title slide ----
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = f"{country} — Quarterly Performance"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = (
                f"Reporting through: {_period_to_label(latest_q)}\n"
                f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            )
        if self.logo_path and os.path.exists(self.logo_path):
            self._add_picture(slide, self.logo_path, left=8.5, top=0.2, width=1.5)

        # ---- KPI slide ----
        self._add_kpi_slide(prs, df_country)

        # ---- Country trend (last N quarters) ----
        self._add_trend_slide(prs, df_country, country, start_q, latest_q)

        # ---- Model mix ----
        top_models = self._add_model_mix_slide(prs, df_country)

        # ---- Top-3 model trends ----
        self._add_top_models_slides(prs, df_country, country, top_models, start_q, latest_q)

        # ---- Save deck ----
        out_name = f"{country}_Quarterly_Performance_{_period_to_label(latest_q)}.pptx"
        out_path = os.path.join(self.output_dir, out_name)
        prs.save(out_path)
        return out_path

    # ---------- Slides ----------

    def _add_kpi_slide(self, prs: Presentation, df_c: pd.DataFrame) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content
        slide.shapes.title.text = "Overview KPIs"

        content = slide.shapes.placeholders[1].text_frame
        content.clear()
        content.word_wrap = True

        latest_q = df_c["QuarterPeriod"].max()
        last_df = df_c[df_c["QuarterPeriod"] == latest_q]

        last_units = float(last_df["Sales_Units"].sum())
        if last_df["Revenue"].notna().any():
            last_rev = float(last_df["Revenue"].sum())
        else:
            last_rev = np.nan

        avg_price = self._safe_div(last_rev, last_units) if np.isfinite(last_rev) else float(last_df["Unit_Price"].mean())

        prev_q = latest_q - 1
        prev_units = float(df_c.loc[df_c["QuarterPeriod"] == prev_q, "Sales_Units"].sum())
        mom_growth = (self._safe_div((last_units - prev_units), prev_units) * 100.0) if prev_units else 0.0

        def add_kpi_bullet(label: str, value: str, unit: str = ""):
            p = content.add_paragraph()
            run = p.add_run()
            run.text = f"{label}: "
            run.font.size = Pt(18)

            run_val = p.add_run()
            run_val.text = f"{value}{unit}"
            run_val.font.bold = True
            run_val.font.size = Pt(18)
            run_val.font.color.rgb = RGBColor(0, 0, 0)

        add_kpi_bullet("Latest Quarter", _period_to_label(latest_q))
        add_kpi_bullet("Total Units", f"{last_units:,.0f}")
        add_kpi_bullet("Revenue", f"{last_rev:,.0f}" if np.isfinite(last_rev) else "—")
        add_kpi_bullet("Avg Unit Price", f"{avg_price:,.0f}" if np.isfinite(avg_price) else "—")
        add_kpi_bullet("QoQ Growth (Units)", f"{mom_growth:.1f}", "%")

    def _add_trend_slide(
        self,
        prs: Presentation,
        df_c: pd.DataFrame,
        country: str,
        start_q: pd.Period,
        latest_q: pd.Period,
    ) -> None:
        trend = (
            df_c[(df_c["QuarterPeriod"] >= start_q) & (df_c["QuarterPeriod"] <= latest_q)]
            .groupby("QuarterPeriod", as_index=False)["Sales_Units"].sum()
            .sort_values("QuarterPeriod")
        )
        if trend.empty:
            return

        # X labels as strings
        x = trend["QuarterPeriod"].apply(_period_to_label)
        y = trend["Sales_Units"]

        plt.figure(figsize=(8, 4))
        plt.plot(x, y)
        plt.title(f"Quarterly Sales Trend — {country} (Last {self.last_n_quarters} quarters)")
        plt.xlabel("Quarter")
        plt.ylabel("Sales Units")
        img_trend = os.path.join(self._tmp_dir, f"{country}_trend.png")
        self._fig_save(img_trend)

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title-only
        slide.shapes.title.text = "Sales Trend"
        self._add_picture(slide, img_trend, left=0.7, top=1.5, width=8.5)

    def _add_model_mix_slide(self, prs: Presentation, df_c: pd.DataFrame) -> List[str]:
        model_mix_df = (
            df_c.groupby("Model", as_index=False)["Sales_Units"]
            .sum()
            .sort_values("Sales_Units", ascending=False)
        )
        if model_mix_df.empty:
            return []

        fig, ax = plt.subplots(figsize=(6, 4))
        bars = ax.bar(model_mix_df["Model"], model_mix_df["Sales_Units"])
        for bar in bars:
            height = bar.get_height()
            ax.annotate(
                f"{height:,.0f}",
                xy=(bar.get_x() + bar.get_width() / 2, height),
                xytext=(0, 3),
                textcoords="offset points",
                ha="center",
                va="bottom",
                fontsize=10,
                fontweight="bold",
            )
        ax.set_title("Model Mix – Units", fontsize=14, fontweight="bold")
        ax.set_ylabel("Units")
        plt.xticks(rotation=30, ha="right")

        chart_path = os.path.join(self._tmp_dir, "model_mix.png")
        plt.tight_layout()
        plt.savefig(chart_path, dpi=300)
        plt.close()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Model Mix"
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))

        top_models = model_mix_df.head(3)["Model"].tolist()
        return top_models

    def _add_top_models_slides(
        self,
        prs: Presentation,
        df_c: pd.DataFrame,
        country: str,
        top_models: Iterable[str],
        start_q: pd.Period,
        latest_q: pd.Period,
    ) -> None:
        for m in top_models:
            d = (
                df_c[
                    (df_c["Model"] == m)
                    & (df_c["QuarterPeriod"] >= start_q)
                    & (df_c["QuarterPeriod"] <= latest_q)
                ]
                .groupby("QuarterPeriod", as_index=False)["Sales_Units"].sum()
                .sort_values("QuarterPeriod")
            )
            if d.empty:
                continue

            plt.figure(figsize=(8, 4))
            plt.plot(d["QuarterPeriod"].apply(_period_to_label), d["Sales_Units"])
            plt.title(f"Model Trend — {m} ({country})")
            plt.xlabel("Quarter")
            plt.ylabel("Sales Units")
            img = os.path.join(self._tmp_dir, f"{country}_{m}_trend.png")
            self._fig_save(img)

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Model Trend — {m}"
            self._add_picture(slide, img, left=0.7, top=1.5, width=8.5)

    # ---------- Data prep ----------

    def _read_table(self, path: str, sheet: Optional[str]) -> pd.DataFrame:
        ext = os.path.splitext(path)[1].lower()
        if ext in (".xlsx", ".xls"):
            df = pd.read_excel(path, sheet_name=sheet)
        else:
            df = pd.read_csv(path)
        df.columns = [c.strip().replace("\t", " ") for c in df.columns]
        return df

    def _prepare_dataframe(self, df: pd.DataFrame, column_map: dict | None = None) -> pd.DataFrame:
        """
        column_map (optional):
          {
            "country": "<exact col name>",
            "model": "<exact col name>",
            "quarter": "<exact col name>",   # values like 'Q1-2025'
            "units": "<exact col name>",
            "revenue": "<exact col name or None>",
            "price": "<exact col name or None>"
          }
        """
        def norm(s: str) -> str:
            return s.lower().replace(" ", "").replace("_", "")

        cols = list(df.columns)

        if column_map:
            def _pick(k, required=True):
                v = column_map.get(k)
                if required and (not v or v not in cols):
                    raise KeyError(f"Column map for '{k}' is missing or invalid.")
                return v if v in cols else None

            col_country = _pick("country", required=True)
            col_model   = _pick("model",   required=True)
            col_quarter = _pick("quarter", required=True)
            col_units   = _pick("units",   required=True)
            col_revenue = _pick("revenue", required=False)
            col_price   = _pick("price",   required=False)
        else:
            def find_one(candidates: tuple[str, ...]) -> str:
                for c in cols:
                    if norm(c) in candidates:
                        return c
                raise KeyError(f"Missing required column among: {candidates}")

            def find_optional(candidates: tuple[str, ...]) -> Optional[str]:
                for c in cols:
                    if norm(c) in candidates:
                        return c
                return None

            col_country = find_one(("country",))
            col_model   = find_one(("bikemodel", "model", "bikemode", "modelname"))
            col_quarter = find_one(("quarter", "qtr", "qrtr", "period"))
            col_units   = find_one(("salesunits", "units", "sales", "qty", "quantity"))
            col_revenue = find_optional(("revenue", "revenueinr", "amount", "salesvalue"))
            col_price   = find_optional(("unitprice", "price", "asp", "avgprice", "averageprice"))

        work = df.copy()

        # ---- Quarter parsing ----
        qp = work[col_quarter].apply(_quarter_to_period)
        work = work.loc[qp.notna()].copy()
        work["QuarterPeriod"] = qp[qp.notna()].astype("period[Q]")
        work["Quarter"] = work["QuarterPeriod"].apply(_period_to_label)

        # Canonical columns
        work = work.rename(columns={col_country: "Country", col_model: "Model"})
        work["Sales_Units"] = pd.to_numeric(work[col_units], errors="coerce").fillna(0)

        if col_revenue:
            work["Revenue"] = pd.to_numeric(work[col_revenue], errors="coerce")
        else:
            work["Revenue"] = np.nan

        if col_price:
            work["Unit_Price"] = pd.to_numeric(work[col_price], errors="coerce")
        else:
            work["Unit_Price"] = np.nan

        # Derive Unit_Price if missing but Revenue & Units present
        mask = work["Unit_Price"].isna() & work["Revenue"].notna() & (work["Sales_Units"] > 0)
        work.loc[mask, "Unit_Price"] = work.loc[mask, "Revenue"] / work.loc[mask, "Sales_Units"]

        return work

    # ---------- Utils ----------

    @staticmethod
    def _ensure_dir(p: str) -> None:
        os.makedirs(p, exist_ok=True)

    @staticmethod
    def _safe_div(a, b) -> float:
        try:
            if b is None or (isinstance(b, float) and np.isnan(b)) or b == 0:
                return np.nan
            return float(a) / float(b)
        except Exception:
            return np.nan

    @staticmethod
    def _fig_save(path: str) -> None:
        plt.tight_layout()
        plt.savefig(path, dpi=200, bbox_inches="tight")
        plt.close()

    @staticmethod
    def _add_picture(slide, image_path: str, left: float, top: float, width: float) -> None:
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
